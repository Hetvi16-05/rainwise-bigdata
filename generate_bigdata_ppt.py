from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_big_data_ppt():
    prs = Presentation()
    
    # --- HELPER: ADD SLIDE ---
    def add_slide(title_text, content_points, image_path=None):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        tf.text = content_points[0]
        for point in content_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(1), Inches(4.5), height=Inches(2.5))

    # --- SLIDE 1: TITLE ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "RAINWISE V3: Big Data Engineering"
    subtitle.text = "Distributed Ingestion, Hadoop Storage, and Spark Processing\n(Focus: Infrastructure & Data Lifecycle)"

    # --- SLIDE 2: THE 5 V'S OF RAINWISE ---
    add_slide("The 5 V's of RAINWISE Big Data", [
        "Volume: 2.2 Million production records + Daily Live Logs.",
        "Velocity: Hourly automated ingestion cycles via Cron.",
        "Variety: Structured CSVs, Semi-structured JSON APIs, and Unstructured Logs.",
        "Veracity: Multi-stage Spark cleaning & schema validation.",
        "Value: Real-time flood hazard intelligence and demographic projections."
    ])

    # --- SLIDE 3: HIGH-LEVEL ARCHITECTURE ---
    add_slide("High-Level Pipeline Architecture", [
        "Ingestion: NASA, CWC, and Open-Meteo APIs (Python/Requests).",
        "Staging: Temporary local storage with automated purge (Zero Local Footprint).",
        "Storage: Official Hadoop HDFS (NameNode & DataNode blocks).",
        "Processing: Apache Spark (PySpark) for distributed transformations.",
        "Visualization: RAINWISE Command Center Dashboard."
    ], image_path="/Users/HetviSheth/.gemini/antigravity/brain/1628f95e-ca0d-436c-9c12-950c4301868f/big_data_architecture_diagram_1777136324138.png")

    # --- SLIDE 4: HDFS STORAGE LAYER (LIVE) ---
    add_slide("Storage Layer: Official Hadoop HDFS", [
        "Namespace: Managed by NameNode at /user/HetviSheth/rainwise/.",
        "Status: ACTIVE (Verified via NameNode Web UI - Port 9870).",
        "Data Replication: Default replication factor of 3 for fault tolerance.",
        "Block Management: Live CSV data stored in 128MB HDFS blocks.",
        "Official Bridge: Subprocess integration with 'hadoop fs -put' for real-time sync."
    ], image_path="/Users/HetviSheth/.gemini/antigravity/brain/1628f95e-ca0d-436c-9c12-950c4301868f/hadoop_ui_live_1777137018506.png")

    # --- SLIDE 5: PARALLELISM & DISTRIBUTED PROCESSING ---
    add_slide("Parallelism & Distributed Processing", [
        "Partitioning: Spark RDDs are split into multiple partitions (e.g., 16-32 partitions).",
        "Concurrent Tasks: 16+ tasks run in parallel across CPU cores.",
        "In-Memory Shuffling: Distributed 'Join' and 'GroupBy' operations via DAGs.",
        "Resource Manager: YARN (Yet Another Resource Negotiator) manages cluster capacity.",
        "Scalability: Performance scales linearly with the number of DataNodes/Cores."
    ], image_path="/Users/HetviSheth/.gemini/antigravity/brain/1628f95e-ca0d-436c-9c12-950c4301868f/.system_generated/click_feedback/click_feedback_1777137137985.png")

    # --- SLIDE 6: DATA VERACITY & GOVERNANCE ---
    add_slide("Data Veracity & Governance", [
        "Schema Enforcement: Standardizing formats across 213 city stations.",
        "Anomaly Detection: Spark-based filtering for sensor outliers.",
        "Consistency: Mapping daily rainfall to historical climate baselines.",
        "Audit Trail: Automated pipeline logging for every ingestion cycle.",
        "Result: High-fidelity 'Golden Zone' data for accurate hazard assessment."
    ])

    # --- SLIDE 7: AUTOMATION & ORCHESTRATION ---
    add_slide("Pipeline Automation", [
        "Orchestrator: run_realtime_pipeline.py managing 4 sub-modules.",
        "Scheduling: Cron-based automation (runs every 60 minutes).",
        "Lock System: Safety checks to prevent concurrent pipeline collisions.",
        "Bridge Logic: Python-to-Hadoop communication using standard CLI interfaces.",
        "Lifecycle: Fetch -> Bridge -> Retrain -> Purge (Automated Cycle)."
    ])

    # --- SLIDE 8: SUMMARY ---
    add_slide("Summary & Conclusion", [
        "100% Native Hadoop infrastructure for storage.",
        "Decoupled Storage (HDFS) and Processing (Spark).",
        "Zero Local Footprint for real-time live data.",
        "Scalable to millions of records with distributed fault-tolerance.",
        "Ready for Production Deployment."
    ])

    save_path = "/Users/HetviSheth/rainwise/RAINWISE_BigData_Pipeline.pptx"
    prs.save(save_path)
    print(f"✅ Presentation saved: {save_path}")

if __name__ == "__main__":
    create_big_data_ppt()
